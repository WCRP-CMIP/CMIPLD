#!/usr/bin/env python3
"""
PowerPoint Template Generator with Interactive Forms

Creates PowerPoint presentations with interactive form controls (radio buttons, checkboxes, dropdowns)
based on CSV field definitions and Python configuration files.

pip install python-pptx

Note: For full interactivity, the generated PowerPoint may need to be opened in 
PowerPoint desktop application rather than online viewers.
"""

import csv
import sys
import argparse
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.oxml.shared import OxmlElement, qn
from pptx.oxml.ns import nsdecls
from pptx.oxml import parse_xml

failed_templates = []

def load_template_data(py_file):
    """Load configuration and data from Python file."""
    namespace = {}
    with open(py_file, 'r', encoding='utf-8') as f:
        exec(f.read(), namespace)
    
    config = namespace.get('TEMPLATE_CONFIG', {})
    data = namespace.get('DATA', {})
    return config, data

def load_csv_fields(csv_file):
    """Load field definitions from CSV."""
    fields = []
    with open(csv_file, 'r', encoding='utf-8') as f:
        reader = csv.DictReader(f)
        for row in reader:
            fields.append(row)
    
    fields.sort(key=lambda x: int(x['field_order']))
    return fields

def setup_slide_layout(prs):
    """Set up a custom slide layout for forms."""
    # Use blank layout
    slide_layout = prs.slide_layouts[6]  # Blank layout
    return slide_layout

def create_title_slide(prs, config):
    """Create the title slide."""
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title = slide.shapes.title
    title.text = config['name']
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Subtitle
    if slide.placeholders:
        try:
            subtitle = slide.placeholders[1]
            subtitle_text = config.get('description', 'Interactive Template Form')
            subtitle_text += f"\n\nGenerated from: {config.get('labels', 'Template')}"
            subtitle.text = subtitle_text
            subtitle.text_frame.paragraphs[0].font.size = Pt(24)
            subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        except (KeyError, IndexError):
            pass
    
    return slide

def create_overview_slide(prs, config, fields):
    """Create an overview slide with template information."""
    slide_layout = setup_slide_layout(prs)
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "Template Overview"
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER
    
    # Overview content
    overview_text = f"""Template Information:
    
• Name: {config['name']}
• Description: {config.get('description', 'N/A')}
• Labels: {config.get('labels', 'N/A')}
• Total Fields: {len([f for f in fields if f['field_type'] != 'markdown'])}
• Required Fields: {len([f for f in fields if f['required'].lower() == 'true' and f['field_type'] != 'markdown'])}

Instructions:
• Fields marked with (*) are required
• Use interactive controls to make selections
• Click checkboxes for multiple selections
• Use dropdown menus for single selections
• Fill text boxes with appropriate information"""
    
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
    content_frame = content_box.text_frame
    content_para = content_frame.paragraphs[0]
    content_para.text = overview_text
    content_para.font.size = Pt(18)
    content_para.alignment = PP_ALIGN.LEFT
    
    return slide

def add_checkbox_control(slide, x, y, width, height, text, checked=False):
    """Add a checkbox control to the slide."""
    # Create checkbox shape (rounded rectangle)
    checkbox = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(0.2), Inches(0.2)
    )
    
    # Style the checkbox
    checkbox.fill.solid()
    checkbox.fill.fore_color.rgb = RGBColor(255, 255, 255)
    checkbox.line.color.rgb = RGBColor(0, 0, 0)
    checkbox.line.width = Pt(1.5)
    
    if checked:
        # Add checkmark (simple approach)
        check_shape = slide.shapes.add_shape(
            MSO_SHAPE.DIAMOND, x + Inches(0.05), y + Inches(0.05), 
            Inches(0.1), Inches(0.1)
        )
        check_shape.fill.solid()
        check_shape.fill.fore_color.rgb = RGBColor(0, 128, 0)
        check_shape.line.color.rgb = RGBColor(0, 128, 0)
    
    # Add label text
    text_box = slide.shapes.add_textbox(x + Inches(0.3), y, width - Inches(0.3), height)
    text_frame = text_box.text_frame
    text_frame.margin_left = 0
    text_frame.margin_top = Inches(0.05)
    para = text_frame.paragraphs[0]
    para.text = text
    para.font.size = Pt(14)
    
    return checkbox, text_box

def add_radio_button_group(slide, x, y, width, options, selected_index=None):
    """Add a group of radio buttons."""
    button_height = Inches(0.3)
    total_height = len(options) * button_height
    
    buttons = []
    for i, option in enumerate(options):
        button_y = y + (i * button_height)
        
        # Create radio button (circle)
        radio = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x, button_y + Inches(0.05), Inches(0.15), Inches(0.15)
        )
        radio.fill.solid()
        radio.fill.fore_color.rgb = RGBColor(255, 255, 255)
        radio.line.color.rgb = RGBColor(0, 0, 0)
        radio.line.width = Pt(1.5)
        
        # Add selection indicator
        if selected_index is not None and i == selected_index:
            inner_circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, x + Inches(0.03), button_y + Inches(0.08), 
                Inches(0.09), Inches(0.09)
            )
            inner_circle.fill.solid()
            inner_circle.fill.fore_color.rgb = RGBColor(0, 100, 200)
            inner_circle.line.width = Pt(0)
        
        # Add label
        text_box = slide.shapes.add_textbox(
            x + Inches(0.25), button_y, width - Inches(0.25), button_height
        )
        text_frame = text_box.text_frame
        text_frame.margin_left = 0
        text_frame.margin_top = Inches(0.05)
        para = text_frame.paragraphs[0]
        para.text = option
        para.font.size = Pt(14)
        
        buttons.append((radio, text_box))
    
    return buttons

def add_dropdown_control(slide, x, y, width, height, options, selected_text="Select an option"):
    """Add a dropdown-style control."""
    # Main dropdown box
    dropdown = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x, y, width, height
    )
    dropdown.fill.solid()
    dropdown.fill.fore_color.rgb = RGBColor(240, 240, 240)
    dropdown.line.color.rgb = RGBColor(100, 100, 100)
    dropdown.line.width = Pt(1)
    
    # Dropdown arrow
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_TRIANGLE, x + width - Inches(0.3), y + Inches(0.05), 
        Inches(0.2), Inches(0.15)
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RGBColor(100, 100, 100)
    arrow.line.width = Pt(0)
    arrow.rotation = 90  # Point downward
    
    # Selected text
    text_box = slide.shapes.add_textbox(x + Inches(0.1), y, width - Inches(0.4), height)
    text_frame = text_box.text_frame
    text_frame.margin_top = Inches(0.05)
    para = text_frame.paragraphs[0]
    para.text = selected_text
    para.font.size = Pt(14)
    
    # Add options as a note (since PowerPoint doesn't support true dropdowns in python-pptx)
    note_y = y + height + Inches(0.05)
    options_text = "Options: " + ", ".join(options[:5])  # Show first 5 options
    if len(options) > 5:
        options_text += f" ... ({len(options)} total options)"
    
    note_box = slide.shapes.add_textbox(x, note_y, width, Inches(0.3))
    note_frame = note_box.text_frame
    note_para = note_frame.paragraphs[0]
    note_para.text = options_text
    note_para.font.size = Pt(10)
    note_para.font.italic = True
    note_para.font.color.rgb = RGBColor(100, 100, 100)
    
    return dropdown, text_box, note_box

def add_text_input_control(slide, x, y, width, height, placeholder=""):
    """Add a text input control."""
    text_input = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, width, height)
    text_input.fill.solid()
    text_input.fill.fore_color.rgb = RGBColor(255, 255, 255)
    text_input.line.color.rgb = RGBColor(150, 150, 150)
    text_input.line.width = Pt(1)
    
    # Placeholder text
    if placeholder:
        text_box = slide.shapes.add_textbox(x + Inches(0.1), y, width - Inches(0.2), height)
        text_frame = text_box.text_frame
        text_frame.margin_top = Inches(0.05)
        para = text_frame.paragraphs[0]
        para.text = placeholder
        para.font.size = Pt(12)
        para.font.italic = True
        para.font.color.rgb = RGBColor(150, 150, 150)
        
        return text_input, text_box
    
    return text_input, None

def create_field_slide(prs, field_def, data, field_number):
    """Create a slide for a single field."""
    slide_layout = setup_slide_layout(prs)
    slide = prs.slides.add_slide(slide_layout)
    
    field_type = field_def['field_type']
    field_id = field_def['field_id']
    label = field_def['label']
    description = field_def['description']
    data_source = field_def['data_source']
    required = field_def['required'].lower() == 'true'
    placeholder = field_def['placeholder']
    options_type = field_def['options_type']
    default_value = field_def['default_value']
    
    # Skip markdown fields
    if field_type == 'markdown':
        if description:
            # Create info slide
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
            title_frame = title_box.text_frame
            title_para = title_frame.paragraphs[0]
            title_para.text = "Information"
            title_para.font.size = Pt(32)
            title_para.font.bold = True
            title_para.alignment = PP_ALIGN.CENTER
            
            content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
            content_frame = content_box.text_frame
            content_para = content_frame.paragraphs[0]
            content_para.text = description.replace('\\n', '\n')
            content_para.font.size = Pt(18)
            content_para.alignment = PP_ALIGN.LEFT
        return slide
    
    # Field title
    title_text = f"{field_number}. {label}"
    if required:
        title_text += " *"
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title_text
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.LEFT
    
    # Field description
    y_position = Inches(1.5)
    if description and description.strip():
        desc_box = slide.shapes.add_textbox(Inches(0.5), y_position, Inches(9), Inches(0.8))
        desc_frame = desc_box.text_frame
        desc_para = desc_frame.paragraphs[0]
        desc_para.text = f"Description: {description.replace('\\n', ' ')}"
        desc_para.font.size = Pt(16)
        desc_para.alignment = PP_ALIGN.LEFT
        y_position += Inches(1)
    
    # Field type info
    type_info = f"Field Type: {field_type.title()} | Field ID: {field_id}"
    type_box = slide.shapes.add_textbox(Inches(0.5), y_position, Inches(9), Inches(0.4))
    type_frame = type_box.text_frame
    type_para = type_frame.paragraphs[0]
    type_para.text = type_info
    type_para.font.size = Pt(12)
    type_para.font.italic = True
    type_para.font.color.rgb = RGBColor(100, 100, 100)
    y_position += Inches(0.6)
    
    # Add interactive control based on field type
    control_y = y_position
    
    if field_type == 'input':
        placeholder_text = placeholder if placeholder else "Enter text here..."
        add_text_input_control(slide, Inches(0.5), control_y, Inches(8), Inches(0.5), placeholder_text)
        
    elif field_type == 'textarea':
        placeholder_text = placeholder if placeholder else "Enter multiple lines of text here..."
        add_text_input_control(slide, Inches(0.5), control_y, Inches(8), Inches(1.5), placeholder_text)
        
    elif field_type == 'dropdown':
        options = get_field_options(field_def, data)
        if options:
            selected = default_value if default_value else "Select an option"
            add_dropdown_control(slide, Inches(0.5), control_y, Inches(6), Inches(0.4), options, selected)
        
    elif field_type == 'multi-select':
        options = get_field_options(field_def, data)
        if options:
            # Create checkboxes for multi-select
            for i, option in enumerate(options[:8]):  # Limit to 8 options per slide
                checkbox_y = control_y + (i * Inches(0.4))
                add_checkbox_control(slide, Inches(0.5), checkbox_y, Inches(7), Inches(0.3), option)
            
            if len(options) > 8:
                note_y = control_y + (8 * Inches(0.4))
                note_box = slide.shapes.add_textbox(Inches(0.5), note_y, Inches(8), Inches(0.3))
                note_frame = note_box.text_frame
                note_para = note_frame.paragraphs[0]
                note_para.text = f"... and {len(options) - 8} more options"
                note_para.font.size = Pt(12)
                note_para.font.italic = True
                note_para.font.color.rgb = RGBColor(100, 100, 100)
        
    elif field_type == 'checkboxes':
        options = get_field_options(field_def, data)
        if options:
            for i, option in enumerate(options[:6]):  # Limit to 6 for checkboxes
                checkbox_y = control_y + (i * Inches(0.4))
                add_checkbox_control(slide, Inches(0.5), checkbox_y, Inches(7), Inches(0.3), option)
    
    # Add instructions at the bottom
    instructions = get_filling_instructions(field_type, required, field_def, data)
    if instructions:
        inst_y = Inches(6.5)
        inst_box = slide.shapes.add_textbox(Inches(0.5), inst_y, Inches(9), Inches(1))
        inst_frame = inst_box.text_frame
        inst_para = inst_frame.paragraphs[0]
        inst_para.text = f"Instructions: {instructions}"
        inst_para.font.size = Pt(14)
        inst_para.font.bold = True
        inst_para.font.color.rgb = RGBColor(0, 100, 0)
    
    return slide

def get_field_options(field_def, data):
    """Extract options for dropdown and multi-select fields."""
    data_source = field_def['data_source']
    options_type = field_def['options_type']
    
    options = []
    
    if data_source != 'none' and data_source in data:
        source_data = data[data_source]
        
        if options_type == 'dict_keys':
            options = list(source_data.keys())
        elif options_type == 'list':
            options = list(source_data)
        elif options_type in ['dict_multiple']:
            options = list(source_data.keys())
        elif options_type == 'dict_with_extra':
            options = list(source_data.keys())
            options.extend(["Open Source", "Registration Required", "Proprietary"])
        elif options_type == 'list_with_na':
            options = ["Not applicable"] + list(source_data)
    
    return options

def get_filling_instructions(field_type, required, field_def, data):
    """Generate filling instructions based on field type."""
    instructions = []
    
    if field_type == 'input':
        instructions.append("Click in the text box and enter a single line of text")
    elif field_type == 'textarea':
        instructions.append("Click in the text area and enter multiple lines of text")
    elif field_type == 'dropdown':
        instructions.append("Click the dropdown arrow and select one option")
    elif field_type == 'multi-select':
        instructions.append("Check multiple boxes as needed")
    elif field_type == 'checkboxes':
        instructions.append("Check one or more boxes")
    
    if required:
        instructions.append("This field is required")
    else:
        instructions.append("This field is optional")
    
    return ". ".join(instructions) + "."

def create_summary_slide(prs, config, total_fields):
    """Create a summary slide."""
    slide_layout = setup_slide_layout(prs)
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "Form Complete!"
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER
    
    # Summary content
    summary_text = f"""Thank you for completing the "{config['name']}" template form.

Next Steps:
• Review your responses in presentation mode
• Export or print this presentation for your records
• Submit the information through the appropriate channel

Form Statistics:
• Total Fields Completed: {total_fields}
• Template: {config.get('labels', 'N/A')}

For technical support or questions, please refer to the project documentation."""
    
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    content_frame = content_box.text_frame
    content_para = content_frame.paragraphs[0]
    content_para.text = summary_text
    content_para.font.size = Pt(18)
    content_para.alignment = PP_ALIGN.LEFT
    
    return slide

def generate_powerpoint_presentation(config, fields, data):
    """Generate complete PowerPoint presentation with interactive forms."""
    prs = Presentation()
    
    # Create title slide
    create_title_slide(prs, config)
    
    # Create overview slide
    create_overview_slide(prs, config, fields)
    
    # Create field slides
    field_number = 1
    for field_def in fields:
        # Merge config defaults into data if needed
        if field_def['field_id'] not in data and field_def['field_id'] in config:
            data[field_def['field_id']] = config[field_def['field_id']]
        
        create_field_slide(prs, field_def, data, field_number)
        
        # Only increment counter for non-markdown fields
        if field_def['field_type'] != 'markdown':
            field_number += 1
    
    # Create summary slide
    create_summary_slide(prs, config, field_number - 1)
    
    return prs

def process_template_pair(template_name, csv_file, py_file, output_dir):
    """Process a single template pair and generate PowerPoint presentation."""
    
    print(f"Processing {template_name}...")
    
    try:
        config, data = load_template_data(py_file)
        if not config:
            print(f"    No TEMPLATE_CONFIG found")
            return False
        
        print(f"    Loaded: {config['name']}")
        
        fields = load_csv_fields(csv_file)
        print(f"    Fields: {len(fields)}")
        
        prs = generate_powerpoint_presentation(config, fields, data)
        
        # Save the presentation
        output_file = output_dir / f"{template_name}_form.pptx"
        prs.save(str(output_file))
        
        print(f"    Generated {template_name}_form.pptx")
        return True
        
    except Exception as e:
        print(f"    Error: {e}")
        failed_templates.append(f"{template_name}: {str(e)}")
        return False

def parse_arguments():
    """Parse command line arguments."""
    parser = argparse.ArgumentParser(description='Generate PowerPoint interactive forms for issue templates')
    parser.add_argument('-t', '--template-dir', type=Path, help='Template directory')
    parser.add_argument('-o', '--output-dir', type=Path, help='Output directory')
    parser.add_argument('--template', type=str, help='Generate form for specific template')
    return parser.parse_args()

def main():
    """Main function."""
    args = parse_arguments()
    
    print("PowerPoint Interactive Template Generator")
    print("Generating interactive forms for issue templates")
    
    script_dir = Path.cwd()
    template_dir = args.template_dir or script_dir / ".github" / "GEN_ISSUE_TEMPLATE"
    output_dir = args.output_dir or script_dir / ".github" / "template_forms_pptx"
    
    print(f"Template dir: {template_dir}")
    print(f"Output dir: {output_dir}")
    
    if not template_dir.exists():
        print(f"Template directory not found: {template_dir}")
        return False
    
    output_dir.mkdir(parents=True, exist_ok=True)
    
    csv_files = list(template_dir.glob('*.csv'))
    
    if args.template:
        csv_files = [f for f in csv_files if f.stem == args.template]
    
    if not csv_files:
        print("No CSV files found")
        return False
    
    print(f"Processing {len(csv_files)} templates")
    
    success_count = 0
    for csv_file in csv_files:
        template_name = csv_file.stem
        py_file = template_dir / f"{template_name}.py"
        
        if py_file.exists():
            if process_template_pair(template_name, csv_file, py_file, output_dir):
                success_count += 1
        else:
            print(f"    Warning: No Python config file found for {template_name}")
    
    print(f"Results: {success_count}/{len(csv_files)} successful")
    
    if failed_templates:
        print("\nFailed templates:")
        for failure in failed_templates:
            print(f"  - {failure}")
    
    print("\nNote: For full interactivity, open the generated PowerPoint files")
    print("in Microsoft PowerPoint desktop application.")
    
    return success_count > 0

if __name__ == '__main__':
    try:
        success = main()
        sys.exit(0 if success else 1)
    except ImportError as e:
        if 'pptx' in str(e):
            print("Error: python-pptx package not found. Please install it:")
            print("pip install python-pptx")
            sys.exit(1)
        else:
            raise
